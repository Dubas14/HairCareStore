#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Extract product data from catalog files (PDF, XLS, PPTX) for HAIR LAB store.
Outputs JSON files to scripts/seed-data/
"""
import sys
import os
import re
import json

sys.stdout.reconfigure(encoding='utf-8')

import pandas as pd
import fitz  # PyMuPDF

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
CATALOG_DIR = os.path.join(os.path.dirname(BASE_DIR), 'NO_GIT_ONLY_DEV_CATALOGE')
SEED_DIR = os.path.join(BASE_DIR, 'seed-data')
IMG_DIR = os.path.join(SEED_DIR, 'images')

os.makedirs(SEED_DIR, exist_ok=True)
os.makedirs(IMG_DIR, exist_ok=True)

def find_file(keyword, ext=None):
    """Find a file in catalog dir by keyword."""
    for f in os.listdir(CATALOG_DIR):
        if keyword in f and (ext is None or f.endswith(ext)):
            return os.path.join(CATALOG_DIR, f)
    return None

def slugify(text):
    """Create URL-friendly slug from Ukrainian text."""
    transliteration = {
        'а': 'a', 'б': 'b', 'в': 'v', 'г': 'h', 'ґ': 'g', 'д': 'd', 'е': 'e',
        'є': 'ye', 'ж': 'zh', 'з': 'z', 'и': 'y', 'і': 'i', 'ї': 'yi', 'й': 'y',
        'к': 'k', 'л': 'l', 'м': 'm', 'н': 'n', 'о': 'o', 'п': 'p', 'р': 'r',
        'с': 's', 'т': 't', 'у': 'u', 'ф': 'f', 'х': 'kh', 'ц': 'ts', 'ч': 'ch',
        'ш': 'sh', 'щ': 'shch', 'ь': '', 'ю': 'yu', 'я': 'ya', 'ъ': '', 'ы': 'y',
        'э': 'e', "'": '',
    }
    text = text.lower().strip()
    result = []
    for ch in text:
        if ch in transliteration:
            result.append(transliteration[ch])
        elif ch.isascii() and (ch.isalnum() or ch == '-'):
            result.append(ch)
        elif ch in (' ', '_', '.', '/'):
            result.append('-')
    slug = '-'.join(part for part in ''.join(result).split('-') if part)
    return slug[:80]

def extract_volume(name):
    """Extract volume from product name like '250 мл', '500 гр'."""
    m = re.search(r'(\d+)\s*(мл|гр|ml|gr|g)\b', name, re.IGNORECASE)
    if m:
        return f"{m.group(1)} {m.group(2).lower()}"
    m = re.search(r'(\d+)\s*(мл|гр)', name)
    if m:
        return f"{m.group(1)} {m.group(2)}"
    return None


# ═══════════════════════════════════════════════════════════════
# 1. ELGON XLS
# ═══════════════════════════════════════════════════════════════
def parse_elgon():
    print("📦 Parsing Elgon XLS...")
    f = find_file('Elgon', '.xls')
    if not f:
        print("  ⚠ Elgon XLS not found")
        return []

    df = pd.read_excel(f, header=None)
    products = []
    current_categories = []  # stack of category names by indent level

    for i in range(10, len(df)):
        row = df.iloc[i]
        art = str(row[1]).strip() if pd.notna(row[1]) else ''
        code = str(row[2]).strip() if pd.notna(row[2]) else ''
        name_raw = str(row[3]) if pd.notna(row[3]) else ''
        name = name_raw.strip()
        cost_val = row[4] if pd.notna(row[4]) else None
        price_val = row[5] if pd.notna(row[5]) else None

        if not name:
            continue

        indent = len(name_raw) - len(name_raw.lstrip())

        # Category row (no article)
        if not art:
            # Trim stack to current indent level
            current_categories = [(lvl, cat) for lvl, cat in current_categories if lvl < indent]
            current_categories.append((indent, name))
            continue

        # Product row
        try:
            cost = int(float(cost_val)) if cost_val else None
            price = int(float(price_val)) if price_val else None
        except (ValueError, TypeError):
            cost = None
            price = None

        if not price and not cost:
            continue

        # Determine category from stack
        cat_name = ''
        if current_categories:
            # Use the deepest meaningful category
            cat_name = current_categories[-1][1] if current_categories else ''

        products.append({
            'title': name,
            'brand': 'elgon',
            'categoryHint': cat_name,
            'articleCode': art,
            'supplierCode': code,
            'price': price,
            'costPrice': cost,
            'volume': extract_volume(name),
            'inStock': True,
        })

    print(f"  ✅ Elgon: {len(products)} products")
    return products


# ═══════════════════════════════════════════════════════════════
# 2. MOOD PDF
# ═══════════════════════════════════════════════════════════════
def parse_mood():
    print("📦 Parsing MOOD Price PDF...")
    f = find_file('MOOD', '.pdf')
    # Find the price PDF specifically
    for fn in os.listdir(CATALOG_DIR):
        if 'MOOD' in fn and '2026' in fn and fn.endswith('.pdf'):
            f = os.path.join(CATALOG_DIR, fn)
            break

    if not f:
        print("  ⚠ MOOD price PDF not found")
        return []

    doc = fitz.open(f)
    full_text = ''
    for page in doc:
        full_text += page.get_text() + '\n'
    doc.close()

    products = []

    # Parse structured data: look for patterns like
    # SKU number + volume + prices
    lines = full_text.split('\n')
    i = 0
    current_section = ''
    current_product_name = ''

    while i < len(lines):
        line = lines[i].strip()

        # Track section headers (all caps or known sections)
        if line and len(line) > 3 and not any(c.isdigit() for c in line[:3]):
            # Check if this might be a section or product name
            if any(kw in line.upper() for kw in ['DREAM CURLS', 'BODY BUILDER', 'ULTRA CARE', 'KERATIN',
                'INTENSE REPAIR', 'SILVER SPECIFIC', 'DERMA CLEANSING', 'COLOR PROTECT',
                'DAILY', 'CELL FORCE', 'BODYGUARD', 'SUNCARE', 'HAIR BODYGUARD']):
                current_section = line
                i += 1
                continue

        # Look for article codes (6-10 digit numbers that are SKUs)
        sku_match = re.match(r'^(\d{6,10})\s*(?:-\s*)?(\d+\s*(?:мл|гр|ml))?\s*$', line)
        if not sku_match:
            # Try pattern: "SKU - volume"
            sku_match = re.match(r'^(\d{6,10})\s*-\s*(\d+\s*(?:мл|гр|ml))\s*$', line)
        if not sku_match:
            # Try just SKU on its own line
            sku_match = re.match(r'^(\d{6,10})$', line)

        if sku_match:
            sku = sku_match.group(1)
            volume_from_sku = sku_match.group(2) if sku_match.lastindex and sku_match.lastindex >= 2 else None

            # Look backward for product name
            name = ''
            for j in range(i-1, max(i-8, -1), -1):
                prev = lines[j].strip()
                if prev and not re.match(r'^\d+\s*(грн|мл|гр|Vol|%)', prev) and len(prev) > 5:
                    if not re.match(r'^(Ціна|РРЦ|салону)', prev):
                        name = prev
                        break

            # Look forward for prices
            cost_price = None
            retail_price = None
            volume = volume_from_sku

            for j in range(i+1, min(i+10, len(lines))):
                next_line = lines[j].strip()
                # Volume line
                vol_m = re.match(r'^(\d+)\s*(мл|гр|ml)\s*$', next_line)
                if vol_m and not volume:
                    volume = f"{vol_m.group(1)} {vol_m.group(2)}"
                    continue

                # Price line: "NNN грн"
                price_m = re.match(r'^(\d+)\s*грн$', next_line)
                if price_m:
                    val = int(price_m.group(1))
                    if cost_price is None:
                        cost_price = val
                    elif retail_price is None:
                        retail_price = val
                        break
                    continue

                # Check for "Ціна салону" / "РРЦ" labels
                if 'РРЦ' in next_line or 'салону' in next_line:
                    continue

                # Stop at next section or product
                if re.match(r'^(\d{6,10})', next_line):
                    break

            if name and (retail_price or cost_price):
                # For MOOD: first price is cost (salon), second is retail (РРЦ)
                if retail_price and cost_price and retail_price < cost_price:
                    cost_price, retail_price = retail_price, cost_price

                products.append({
                    'title': name,
                    'brand': 'mood',
                    'categoryHint': current_section,
                    'articleCode': sku,
                    'supplierCode': '',
                    'price': retail_price or cost_price,
                    'costPrice': cost_price if retail_price else None,
                    'volume': volume,
                    'inStock': True,
                })

            i += 1
            continue

        i += 1

    # Deduplicate by SKU
    seen = set()
    unique = []
    for p in products:
        if p['articleCode'] not in seen:
            seen.add(p['articleCode'])
            unique.append(p)

    print(f"  ✅ MOOD: {len(unique)} products")
    return unique


# ═══════════════════════════════════════════════════════════════
# 3. NEVITALY PDF
# ═══════════════════════════════════════════════════════════════
def parse_nevitaly():
    print("📦 Parsing Nevitaly catalog PDF...")
    f = None
    for fn in os.listdir(CATALOG_DIR):
        if 'Nevitaly' in fn and fn.endswith('.pdf'):
            f = os.path.join(CATALOG_DIR, fn)
            break

    if not f:
        print("  ⚠ Nevitaly PDF not found")
        return []

    doc = fitz.open(f)
    products = []

    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text()
        lines = text.split('\n')

        current_product_name = ''
        current_section = ''

        i = 0
        while i < len(lines):
            line = lines[i].strip()

            # Track section headers
            for sec in ['NEV COLOR', 'CURL SUBLIME', 'FILLER SUBLIME', 'COLOR SUBLIME',
                       'HYDRA SOURCE', 'SHIMMER', 'PRECIOUS', 'BLONDE SUBLIME', 'BLOND SUBLIME',
                       'STYLING', 'GENTLE', 'SCALP', 'PURIFYING', 'ENERGY', 'SOOTHING',
                       'DETOX', 'AHA', 'TERRAE', 'SYNUOSA']:
                if sec in line.upper():
                    current_section = line
                    break

            # Look for SKU pattern (7 digits)
            sku_m = re.match(r'^(10\d{5})\s*$', line)
            if not sku_m:
                sku_m = re.match(r'^(10\d{5,7})\s*$', line)

            if sku_m:
                sku = sku_m.group(1)

                # Look backward for product name
                name = ''
                for j in range(i-1, max(i-10, -1), -1):
                    prev = lines[j].strip()
                    if prev and len(prev) > 5 and not re.match(r'^[\d\s,\.грн]+$', prev):
                        if not re.match(r'^(Об\'єм|Ціна|мл|грн|pH|рН)', prev):
                            name = prev
                            break

                # Look for volume and prices nearby
                volume = None
                cost_price = None
                retail_price = None

                for j in range(i-5, min(i+10, len(lines))):
                    check = lines[j].strip() if 0 <= j < len(lines) else ''

                    vol_m = re.match(r'^(\d+)$', check)
                    if vol_m and not volume:
                        val = int(vol_m.group(1))
                        if 50 <= val <= 1500:
                            volume = f"{val} мл"

                    price_m = re.match(r'^(\d{3,5})$', check)
                    if price_m:
                        val = int(price_m.group(1))
                        if 200 <= val <= 5000:
                            if cost_price is None:
                                cost_price = val
                            elif retail_price is None:
                                retail_price = val

                if name and (cost_price or retail_price):
                    if retail_price and cost_price and retail_price < cost_price:
                        cost_price, retail_price = retail_price, cost_price

                    products.append({
                        'title': name,
                        'brand': 'nevitaly',
                        'categoryHint': current_section,
                        'articleCode': sku,
                        'supplierCode': '',
                        'price': retail_price or cost_price,
                        'costPrice': cost_price if retail_price else None,
                        'volume': volume,
                        'inStock': True,
                    })

            i += 1

    doc.close()

    # Deduplicate
    seen = set()
    unique = []
    for p in products:
        key = p['articleCode']
        if key not in seen:
            seen.add(key)
            unique.append(p)

    print(f"  ✅ Nevitaly: {len(unique)} products")
    return unique


# ═══════════════════════════════════════════════════════════════
# 4. INEBRYA PDFs
# ═══════════════════════════════════════════════════════════════
def parse_inebrya():
    print("📦 Parsing Inebrya Price PDFs...")
    # Use "магазини" version (higher margin)
    shop_file = None
    salon_file = None
    for fn in os.listdir(CATALOG_DIR):
        if 'Inebrya' in fn and fn.endswith('.pdf'):
            path = os.path.join(CATALOG_DIR, fn)
            if 'магазин' in fn.lower():
                shop_file = path
            elif 'салон' in fn.lower():
                salon_file = path

    if not shop_file:
        print("  ⚠ Inebrya shop PDF not found")
        return []

    doc = fitz.open(shop_file)
    products = []
    current_section = ''

    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text()
        lines = text.split('\n')

        i = 0
        while i < len(lines):
            line = lines[i].strip()

            # Track sections
            for sec in ['НОВИНКИ', 'ФАРБА', 'ОКИСНИК', 'ОСВІТЛЕННЯ', 'BLONDESSE',
                       'ICE CREAM', 'STYLE-IN', 'SHECARE', 'COLOR PERFECT',
                       'KARYN', 'HAIR LIFT', 'Перманентна', 'Деміперманентна']:
                if sec in line:
                    current_section = line
                    break

            # SKU pattern for Inebrya: 7 digits
            sku_m = re.match(r'^(10\d{5})\s*$', line)
            if not sku_m:
                sku_m = re.match(r'^(\d{7})\s*$', line)

            if sku_m:
                sku = sku_m.group(1)

                # Look forward for: product name, volume, prices
                name = ''
                volume = None
                shop_price = None
                retail_price = None

                for j in range(i+1, min(i+12, len(lines))):
                    next_l = lines[j].strip()

                    # Product name: multi-word line with Ukrainian text
                    if not name and len(next_l) > 10 and re.search(r'[а-яА-ЯіІїЇєЄґҐ]', next_l):
                        if not re.match(r'^[\d\s,\.грн]+$', next_l):
                            name = next_l
                            # Check next line too for continuation
                            if j+1 < len(lines):
                                cont = lines[j+1].strip()
                                if cont and re.search(r'[а-яА-ЯіІїЇєЄ]', cont) and not re.match(r'^\d+\s*(мл|гр|грн)', cont):
                                    if len(cont) > 3 and not re.match(r'^(Ціна|РРЦ)', cont):
                                        name += ' ' + cont
                            continue

                    # Volume
                    vol_m = re.match(r'^(\d+)\s*(мл|гр)\s*$', next_l)
                    if vol_m:
                        volume = f"{vol_m.group(1)} {vol_m.group(2)}"
                        continue

                    # Price: "NNN грн"
                    price_m = re.match(r'^(\d+)\s*грн$', next_l)
                    if price_m:
                        val = int(price_m.group(1))
                        if shop_price is None:
                            shop_price = val
                        elif retail_price is None:
                            retail_price = val
                            break

                    # Next SKU means end of current
                    if re.match(r'^\d{7}$', next_l):
                        break

                if name and (shop_price or retail_price):
                    # Shop price = our costPrice, retail = selling price
                    if retail_price and shop_price and retail_price < shop_price:
                        shop_price, retail_price = retail_price, shop_price

                    products.append({
                        'title': name,
                        'brand': 'inebrya',
                        'categoryHint': current_section,
                        'articleCode': sku,
                        'supplierCode': '',
                        'price': retail_price or shop_price,
                        'costPrice': shop_price if retail_price else None,
                        'volume': volume,
                        'inStock': True,
                    })

            i += 1

    doc.close()

    # Deduplicate
    seen = set()
    unique = []
    for p in products:
        if p['articleCode'] not in seen:
            seen.add(p['articleCode'])
            unique.append(p)

    print(f"  ✅ Inebrya: {len(unique)} products")
    return unique


# ═══════════════════════════════════════════════════════════════
# 5. EXTRACT IMAGES from PDFs
# ═══════════════════════════════════════════════════════════════
def extract_images():
    print("🖼️  Extracting images from catalogs...")
    manifest = []
    brands_map = {
        'Elgon_ua': 'elgon', 'Elgon_Scalpcare': 'elgon', 'AFFIXX': 'elgon',
        'YES ESSENTIAL': 'elgon', 'MOOD': 'mood', 'DREAM CURLS': 'mood',
        'BODY BUILDER': 'mood', 'KERATIN': 'mood', 'SUNCARE': 'mood',
        'BODYGUARD': 'mood', 'Nevitaly': 'nevitaly', 'Inebrya': 'inebrya',
    }

    for fn in os.listdir(CATALOG_DIR):
        filepath = os.path.join(CATALOG_DIR, fn)
        brand = 'unknown'
        for key, val in brands_map.items():
            if key.lower() in fn.lower():
                brand = val
                break

        brand_dir = os.path.join(IMG_DIR, brand)
        os.makedirs(brand_dir, exist_ok=True)

        if fn.endswith('.pdf'):
            try:
                doc = fitz.open(filepath)
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    images = page.get_images(full=True)
                    for img_idx, img in enumerate(images):
                        xref = img[0]
                        try:
                            pix = fitz.Pixmap(doc, xref)
                            if pix.width < 150 or pix.height < 150:
                                continue
                            if pix.n - pix.alpha > 3:  # CMYK
                                pix = fitz.Pixmap(fitz.csRGB, pix)

                            img_name = f"{slugify(fn.rsplit('.', 1)[0])}_p{page_num+1}_img{img_idx+1}.png"
                            img_path = os.path.join(brand_dir, img_name)
                            pix.save(img_path)
                            manifest.append({
                                'file': f"images/{brand}/{img_name}",
                                'brand': brand,
                                'source': fn,
                                'page': page_num + 1,
                                'width': pix.width,
                                'height': pix.height,
                            })
                        except Exception as e:
                            pass
                doc.close()
            except Exception as e:
                print(f"  ⚠ Error processing {fn}: {e}")

        elif fn.endswith('.pptx'):
            try:
                from pptx import Presentation
                from pptx.util import Emu
                prs = Presentation(filepath)
                for slide_idx, slide in enumerate(prs.slides):
                    for shape_idx, shape in enumerate(slide.shapes):
                        if shape.shape_type == 13:  # Picture
                            image = shape.image
                            w = shape.width / 914400 * 96 if shape.width else 0
                            h = shape.height / 914400 * 96 if shape.height else 0
                            if w < 150 or h < 150:
                                continue
                            ext = image.content_type.split('/')[-1]
                            if ext == 'jpeg':
                                ext = 'jpg'
                            img_name = f"{slugify(fn.rsplit('.', 1)[0])}_s{slide_idx+1}_img{shape_idx+1}.{ext}"
                            img_path = os.path.join(brand_dir, img_name)
                            with open(img_path, 'wb') as fout:
                                fout.write(image.blob)
                            manifest.append({
                                'file': f"images/{brand}/{img_name}",
                                'brand': brand,
                                'source': fn,
                                'page': slide_idx + 1,
                                'width': int(w),
                                'height': int(h),
                            })
            except Exception as e:
                print(f"  ⚠ Error processing {fn}: {e}")

    print(f"  ✅ Extracted {len(manifest)} images")
    return manifest


# ═══════════════════════════════════════════════════════════════
# 6. CATEGORIES
# ═══════════════════════════════════════════════════════════════
def generate_categories():
    print("📁 Generating categories...")
    categories = [
        {"name": "Фарба для волосся", "slug": "farba-dlya-volossya", "order": 1, "children": [
            {"name": "Перманентна фарба", "slug": "permanentna-farba", "order": 1},
            {"name": "Безаміачна фарба", "slug": "bezamіachna-farba", "order": 2},
            {"name": "Деміперманентна фарба", "slug": "demipermanentna-farba", "order": 3},
            {"name": "Тонуючі засоби", "slug": "tonuyuchi-zasoby", "order": 4},
        ]},
        {"name": "Окислювачі", "slug": "okyslyuvachi", "order": 2, "children": []},
        {"name": "Освітлення та знебарвлення", "slug": "osvitlennya-ta-znebarvlennya", "order": 3, "children": []},
        {"name": "Догляд за волоссям", "slug": "doglyad-za-volossynam", "order": 4, "children": [
            {"name": "Шампуні", "slug": "shampuni", "order": 1},
            {"name": "Маски та бальзами", "slug": "masky-ta-balzamy", "order": 2},
            {"name": "Кондиціонери", "slug": "kondytsionery", "order": 3},
            {"name": "Незмивні засоби", "slug": "nezmyvni-zasoby", "order": 4},
            {"name": "Олії та сироватки", "slug": "oliyi-ta-syrovatky", "order": 5},
            {"name": "Спреї", "slug": "spreyi", "order": 6},
        ]},
        {"name": "Стайлінг", "slug": "staylinh", "order": 5, "children": [
            {"name": "Лаки", "slug": "laky", "order": 1},
            {"name": "Муси та піни", "slug": "musy-ta-piny", "order": 2},
            {"name": "Пасти та воски", "slug": "pasty-ta-vosky", "order": 3},
            {"name": "Термозахист", "slug": "termozakhyst", "order": 4},
        ]},
        {"name": "Лікування волосся", "slug": "likuvannya-volossya", "order": 6, "children": [
            {"name": "Ампули та концентрати", "slug": "ampuly-ta-kontsentraty", "order": 1},
            {"name": "Кератинові засоби", "slug": "keratynovi-zasoby", "order": 2},
            {"name": "Ламінування", "slug": "laminuvannya", "order": 3},
        ]},
        {"name": "Скальп-догляд", "slug": "skalp-doglyad", "order": 7, "children": [
            {"name": "Проти лупи", "slug": "proty-lupy", "order": 1},
            {"name": "Проти випадіння", "slug": "proty-vypadinnya", "order": 2},
            {"name": "Детокс", "slug": "detoks", "order": 3},
        ]},
        {"name": "Для чоловіків", "slug": "dlya-cholovikiv", "order": 8, "children": []},
        {"name": "Сонцезахист", "slug": "sontsezakhyst", "order": 9, "children": []},
        {"name": "Хімічна завивка", "slug": "khimichna-zavyvka", "order": 10, "children": []},
        {"name": "Набори", "slug": "nabory", "order": 11, "children": []},
        {"name": "Кольорові маски", "slug": "kolorovi-masky", "order": 12, "children": []},
        {"name": "Догляд за кольором", "slug": "doglyad-za-kolorom", "order": 13, "children": []},
    ]
    print(f"  ✅ {sum(1 + len(c.get('children', [])) for c in categories)} categories")
    return categories


# ═══════════════════════════════════════════════════════════════
# 7. BRANDS
# ═══════════════════════════════════════════════════════════════
def generate_brands():
    print("🏷️  Generating brands...")
    return [
        {
            "name": "Elgon",
            "slug": "elgon",
            "shortDescription": "Італійський професійний бренд для догляду за волоссям з 1953 року. Інноваційні формули на основі натуральних компонентів.",
            "countryOfOrigin": "Італія",
            "foundedYear": 1953,
            "website": "https://www.elgoncosmetic.com",
            "accentColor": "#2E86AB",
            "benefits": [
                {"icon": "🇮🇹", "title": "Made in Italy", "description": "Виробництво в Мілані з 1953 року"},
                {"icon": "🌿", "title": "Натуральні формули", "description": "Використання рослинних екстрактів та олій"},
                {"icon": "💎", "title": "Професійна якість", "description": "Продукти для салонів та домашнього використання"},
                {"icon": "🔬", "title": "Інноваційні технології", "description": "Постійне впровадження наукових розробок"},
            ],
        },
        {
            "name": "MOOD",
            "slug": "mood",
            "shortDescription": "Професійний італійський бренд з інноваційними формулами для фарбування та догляду за волоссям.",
            "countryOfOrigin": "Італія",
            "foundedYear": 2000,
            "website": "https://www.moodprofessional.com",
            "accentColor": "#E63946",
            "benefits": [
                {"icon": "🎨", "title": "104 відтінки фарби", "description": "Широка палітра кольорів для будь-яких потреб"},
                {"icon": "🧬", "title": "Пептидні формули", "description": "Інноваційний догляд з пептидами для відновлення"},
                {"icon": "🌱", "title": "Веганські формули", "description": "Серія KERATIN без SLES та солі"},
                {"icon": "☀️", "title": "Сонцезахист", "description": "Спеціальна лінія SUNCARE для захисту волосся"},
            ],
        },
        {
            "name": "Nevitaly",
            "slug": "nevitaly",
            "shortDescription": "Італійський бренд, що поєднує досвід, інноваційні наукові розробки та силу рослинної терапії для здоров'я волосся.",
            "countryOfOrigin": "Італія",
            "foundedYear": 2005,
            "website": "https://nevitaly.com.ua",
            "accentColor": "#588157",
            "benefits": [
                {"icon": "🌿", "title": "Фітотерапія", "description": "Формули на основі рослинних екстрактів та ефірних олій"},
                {"icon": "🔬", "title": "Трихологія", "description": "Професійна лінійка трихологічних засобів для шкіри голови"},
                {"icon": "💆", "title": "Комплексний догляд", "description": "Від фарбування до глибокого відновлення волосся"},
                {"icon": "🇮🇹", "title": "Made in Italy", "description": "Італійська якість та традиції"},
            ],
        },
        {
            "name": "Inebrya",
            "slug": "inebrya",
            "shortDescription": "Італійський бренд професійної косметики для волосся з олією насіння льону та алое вера.",
            "countryOfOrigin": "Італія",
            "foundedYear": 2003,
            "website": "https://inebrya.com.ua",
            "accentColor": "#9B5DE5",
            "benefits": [
                {"icon": "💜", "title": "118 відтінків", "description": "Одна з найширших палітр фарб на ринку"},
                {"icon": "🌿", "title": "Натуральні олії", "description": "Олія насіння льону та алое вера у кожному продукті"},
                {"icon": "💎", "title": "Кератинова серія", "description": "Рослинний кератин та мікрокристали сапфіру"},
                {"icon": "🔬", "title": "ICE CREAM серія", "description": "Інноваційні засоби для глибокого відновлення"},
            ],
        },
    ]


# ═══════════════════════════════════════════════════════════════
# 8. BLOG POSTS
# ═══════════════════════════════════════════════════════════════
def generate_blog_posts():
    print("📝 Generating blog posts...")
    posts = [
        {
            "title": "Як надати об'єм тонкому волоссю: секрети MOOD Body Builder",
            "slug": "yak-nadaty-obyem-tonkomu-volossyu",
            "brand": "mood",
            "excerpt": "Професійна лінійка для ущільнення та надання об'єму тонкому й ослабленому волоссю. Збагачена пептидом Pisum Sativum.",
            "tags": ["об'єм", "тонке волосся", "MOOD", "Body Builder"],
            "content": "Тонке та ослаблене волосся потребує особливого догляду. Лінійка MOOD Body Builder створена спеціально для вирішення цієї проблеми. Формула збагачена пептидом Pisum Sativum, який впливає на діаметр волосся, роблячи його візуально густішим і щільнішим без обтяження.\n\nОсновні продукти серії:\n- Body Builder Densifying Shampoo — шампунь для ущільнення\n- Body Builder Densifying Filler — філер для миттєвого ефекту\n\nРекомендований протокол:\n1. Нанесіть шампунь на вологе волосся, масажуйте 2-3 хвилини\n2. Змийте та нанесіть філер на підсушені рушником кінчики\n3. Висушіть феном від коренів для максимального об'єму",
        },
        {
            "title": "Догляд за кучерявим волоссям з MOOD Dream Curls",
            "slug": "doglyad-za-kucheryavym-volossynam-dream-curls",
            "brand": "mood",
            "excerpt": "Професійна лінійка для підкреслення форми локонів. Збагачена гіперферментованою рисовою водою.",
            "tags": ["кучері", "кучеряве волосся", "MOOD", "Dream Curls"],
            "content": "Кучеряве та хвилясте волосся має свої унікальні потреби. Серія MOOD Dream Curls створена для підкреслення форми локонів без обтяження.\n\nГоловний активний компонент — гіперферментована рисова вода, фітокомплекс, отриманий шляхом гіперферментації. Формула сприяє підтриманню оптимального рівня зволоження, еластичності та чіткої форми локонів.\n\nПродукти серії:\n- Dream Curls Shampoo — зволожуючий шампунь\n- Dream Curls Mask — зволожуюча маска\n- Dream Curls Leave In — спрей-кондиціонер\n\nСекрет ідеальних кучерів: після миття не сушіть волосся рушником — натисніть і видавіть воду, нанесіть Leave In та дозвольте висохнути природним шляхом.",
        },
        {
            "title": "Кератиновий догляд: відновлення та зміцнення волосся",
            "slug": "keratynovyy-doglyad-vidnovlennya-zmitsnennya",
            "brand": "mood",
            "excerpt": "Серія MOOD Keratin на основі рослинного кератину та мультимінерального комплексу для відновлення та зміцнення.",
            "tags": ["кератин", "відновлення", "MOOD", "Keratin"],
            "content": "Волосся, яке піддається хімічній обробці, кератиновому випрямленню та моделюванню, потребує посиленого догляду. Серія MOOD KERATIN — це рішення на основі рослинного кератину та мультимінерального комплексу.\n\nСклад мультимінерального комплексу: залізо, мідь, магній, кремній, цинк. Ці мінерали працюють синергетично для відновлення, зволоження та зміцнення волосся.\n\nВеганська формула без SLES і солі — ідеально підходить для волосся після кератинового випрямлення.\n\nРезультат: захист від ламкості, запобігання появі посічених кінчиків, відновлення структури.",
        },
        {
            "title": "Захист волосся та шкіри влітку: лінійка MOOD Suncare",
            "slug": "zakhyst-volossya-ta-shkiry-vlitku-suncare",
            "brand": "mood",
            "excerpt": "Як захистити волосся від сонця, солоної води та хлору. Професійні рекомендації з лінійкою MOOD Suncare.",
            "tags": ["сонцезахист", "літній догляд", "MOOD", "Suncare"],
            "content": "Влітку волосся потребує особливого захисту від UV-випромінювання, солоної морської води та хлору в басейнах.\n\nЛінійка MOOD SUNCARE розроблена спеціально для літнього періоду і забезпечує комплексний захист волосся та шкіри.\n\nПоради для літнього догляду:\n1. Нанесіть захисний засіб перед виходом на сонце\n2. Після купання обов'язково промийте волосся чистою водою\n3. Використовуйте зволожуючу маску 2-3 рази на тиждень\n4. Уникайте термоприладів — дозвольте волоссю висохнути природним шляхом",
        },
        {
            "title": "Професійний стайлінг з Elgon AFFIXX: повний гід",
            "slug": "profesiynyy-staylinh-elgon-affixx",
            "brand": "elgon",
            "excerpt": "Огляд лінійки стайлінгу Elgon AFFIXX — від термозахисту до фінішного блиску. 16 продуктів для будь-якої зачіски.",
            "tags": ["стайлінг", "Elgon", "AFFIXX", "укладання"],
            "content": "Лінійка Elgon AFFIXX — це повний арсенал стайлінгових засобів для професіоналів та домашнього використання.\n\nОсновні продукти:\n- AFFIXX 11 Straight Look — термозахисний спрей для ідеально гладкого волосся\n- AFFIXX 22 Quick Dry — прискорювач сушіння волосся\n- AFFIXX 42 Volume Pump Mousse — мус для створення об'єму\n- AFFIXX 44 Flex Hold Spray Wax — спрей-віск еластичної фіксації\n- AFFIXX 55 Pack Oil — відновлююча олія для волосся\n- AFFIXX 60 Flex Hold Eco Spray — екологічний лак\n- AFFIXX 67 Hair Lift — пудра для прикореневого об'єму\n- AFFIXX 83 Curls Creator — крем для формування локонів\n- AFFIXX 100 Rasta Gum — текстуруюча гума\n- AFFIXX 101 Fix It — лак надсильної фіксації\n\nЛайфхак: для максимального об'єму нанесіть пудру Hair Lift на корені та підсушіть феном, тримаючи голову донизу.",
        },
        {
            "title": "Протоколи догляду за шкірою голови з Elgon Primaria",
            "slug": "protokoly-doglyadu-za-shkiroyu-holovy-elgon",
            "brand": "elgon",
            "excerpt": "Як обрати правильний догляд для шкіри голови: від лупи до випадіння. Серія Elgon Primaria.",
            "tags": ["скальп", "шкіра голови", "Elgon", "Primaria", "лупа", "випадіння"],
            "content": "Здоров'я волосся починається зі здоров'я шкіри голови. Серія Elgon Primaria пропонує комплексний підхід до вирішення проблем скальпу.\n\nПротокол проти лупи:\n1. Purifying Shampoo — шампунь з цинком\n2. Purifying Lotion — лосьйон проти лупи з цинком\n\nПротокол проти випадіння:\n1. Stimulating Shampoo — зміцнюючий шампунь\n2. Anti Hairloss Treatment — лосьйон в ампулах\n\nПротокол для жирної шкіри голови:\n1. Rebalancing Shampoo — шампунь з глиною\n2. Rebalancing Deep Cleansing — засіб глибокого очищення\n\nДля щоденного використання:\n- Biodaily Shampoo — делікатний щоденний шампунь, підходить для всіх типів",
        },
        {
            "title": "Нова серія Elgon YES Essential: краса та догляд",
            "slug": "nova-seriya-elgon-yes-essential",
            "brand": "elgon",
            "excerpt": "Огляд нових ліній YES: Nourish, Hydra, Curls, Shine, Smooth та Daily. Догляд для кожного типу волосся.",
            "tags": ["Elgon", "YES Essential", "догляд", "новинки"],
            "content": "Серія Elgon YES Essential — це нове покоління засобів для професійного догляду. Кожна лінія спрямована на конкретний тип волосся.\n\nYES Nourish — живлення сухого та пошкодженого волосся:\n- Power Source Shampoo, Hyper Nutri Mask, Wonder Nutri Oil, Miracle Night&Day Serum\n\nYES Hydra — зволоження:\n- Beauty Shampoo, Beauty Conditioner\n\nYES Curls — для кучерявого волосся:\n- Hydra Shampoo, Hydra Mask, Nutri Mask, Hydra Spray, Memory Cream, Gentle Shampoo\n\nYES Shine — для блиску:\n- Sparkle Shampoo, Extra Glow Mask, Crystal Water\n\nYES Smooth — для розгладження:\n- Super Control Shampoo, So Sleek Conditioner, Liss Forever Mask, Magic-Coat Spray\n\nYES Daily — щоденний догляд:\n- Everyday Shampoo, Day-By-Day Hydra Mist, No-stress Dry Shampoo",
        },
        {
            "title": "Огляд професійного бренду Nevitaly: від фарби до трихології",
            "slug": "ohlyad-profesiynoho-brendu-nevitaly",
            "brand": "nevitaly",
            "excerpt": "Nevitaly — італійський бренд, що поєднує інноваційні наукові розробки та силу рослинної терапії для здоров'я волосся.",
            "tags": ["Nevitaly", "огляд бренду", "трихологія", "фарба"],
            "content": "Nevitaly — це більше, ніж професійна косметика для волосся. Це італійський бренд, що поєднує багаторічний досвід, інноваційні наукові розробки та силу рослинної терапії.\n\nОсновні лінійки:\n\n🎨 NEV COLOR — професійна фарба:\n- Floressence (безаміачна) — 7 тонів з екстрактом лотосу\n- Niu_Tech (аміачна) — 9 тонів з гранатовим соком\n- Shine_On (деміперманентна) — для оновлення кольору\n\n💆 Серії догляду:\n- Curl Sublime — для кучерявого волосся (олія бабасу)\n- Filler Sublime — для об'єму тонкого волосся\n- Color Sublime — захист фарбованого волосся\n- Hydra Source — зволоження сухого волосся\n\n🔬 Трихологічна лінія:\n- Gentle Micellar Cleanser, AHA Peeling, Detox Peeling\n- Purifying Cleanser та Lotion — проти лупи\n- Scalp Balance Cleanser та Lotion — відновлення балансу\n- Energy Scalp Cleanser та Lotion — проти випадіння",
        },
    ]
    print(f"  ✅ {len(posts)} blog posts")
    return posts


# ═══════════════════════════════════════════════════════════════
# CATEGORY ASSIGNMENT LOGIC
# ═══════════════════════════════════════════════════════════════
def assign_category(product):
    """Assign category slug based on product title and categoryHint."""
    title = product['title'].lower()
    hint = product.get('categoryHint', '').lower()

    # Hair dye / color
    if any(w in title for w in ['фарба', 'крем-фарба', 'color', 'colour']):
        if any(w in title for w in ['безаміачн', 'bionic']):
            return 'bezamіachna-farba'
        if any(w in title for w in ['деміперманент', 'demi']):
            return 'demipermanentna-farba'
        return 'permanentna-farba'

    # Toning
    if any(w in title for w in ['тонуюч', 'i-care', 'i-light', 'тонер', 'tonalight', 'пігмент прямої']):
        return 'tonuyuchi-zasoby'

    # Oxidants
    if any(w in title for w in ['окисл', 'оксид', 'окисник', 'активатор', 'activator', 'oxidant', 'oxydant']):
        return 'okyslyuvachi'

    # Bleach/lightening
    if any(w in title for w in ['пудра', 'знебарвл', 'освітл', 'bleach', 'blonde', 'lightener']):
        return 'osvitlennya-ta-znebarvlennya'

    # Color masks
    if any(w in hint for w in ['кольоров', 'terrae']):
        return 'kolorovi-masky'

    # Perm / waving
    if any(w in title for w in ['завивк', 'waving', 'fixing lotion', 'біозавівк']):
        return 'khimichna-zavyvka'

    # Sets
    if any(w in title for w in ['набір', 'gift box', 'kit']):
        return 'nabory'

    # Men
    if any(w in title for w in ['man ', 'man,', 'для чоловік', 'elgon man']):
        return 'dlya-cholovikiv'
    if 'для чоловік' in hint:
        return 'dlya-cholovikiv'

    # Suncare
    if any(w in title for w in ['сонцезахис', 'suncare', 'aftersun', 'sun ']):
        return 'sontsezakhyst'
    if 'suncare' in hint.lower():
        return 'sontsezakhyst'

    # Scalp care
    if any(w in title for w in ['проти лупи', 'purifying', 'purif']):
        return 'proty-lupy'
    if any(w in title for w in ['проти випад', 'anti hairloss', 'stimulat', 'випадіння', 'anti-hairloss', 'scalp awake']):
        return 'proty-vypadinnya'
    if any(w in title for w in ['детокс', 'detox', 'пілінг', 'peeling', 'глибокого очищення шкіри', 'rebalancing', 'deep clean']):
        return 'detoks'
    if any(w in title for w in ['скальп', 'scalp', 'шкіри голови', 'шкіри гол']):
        return 'skalp-doglyad'

    # Color care
    if any(w in title for w in ['colorcare', 'color protect', 'color sublime', 'за кольором', 'silver shamp', 'silver cond', 'anti-red', 'anti-yellow', 'фіолетовими пігмент', 'нейтралізац']):
        return 'doglyad-za-kolorom'
    if 'silver' in hint.lower() or 'anti-red' in hint.lower() or 'color' in hint.lower():
        if not any(w in title for w in ['фарба', 'color,', 'крем-фарба']):
            return 'doglyad-za-kolorom'

    # Keratin
    if any(w in title for w in ['кератин', 'keratin']):
        return 'keratynovi-zasoby'

    # Lamination
    if any(w in title for w in ['ламінуванн', 'lamination']):
        return 'laminuvannya'

    # Ampoules
    if any(w in title for w in ['ампул', 'лосьйон', 'lotion', 'концентрат', 'treatment']):
        if 'випадіння' in title or 'hairloss' in title:
            return 'proty-vypadinnya'
        if 'лупи' in title or 'purif' in title:
            return 'proty-lupy'
        return 'ampuly-ta-kontsentraty'

    # Styling products
    if any(w in title for w in ['лак ', 'лак,', 'hairspray', 'hair spray', 'fix it', 'eco spray', 'total fix', 'logic style']):
        return 'laky'
    if any(w in title for w in ['мус ', 'мус,', 'mousse', 'піна', 'foam']):
        return 'musy-ta-piny'
    if any(w in title for w in ['паста', 'paste', 'віск', 'wax', 'гума', 'gum', 'гель', 'gel', 'пудра для об', 'hair lift', 'volumizing powder']):
        if 'знебарвл' not in title:
            return 'pasty-ta-vosky'
    if any(w in title for w in ['термозахис', 'thermo', 'heat defend', 'straight look']):
        return 'termozakhyst'
    if any(w in hint for w in ['стайлінг', 'styling', 'affixx', 'bodyguard', 'style-in']):
        return 'staylinh'

    # Care products
    if any(w in title for w in ['шампунь', 'shampoo', 'cleanser', 'cleancer']):
        return 'shampuni'
    if any(w in title for w in ['маска', 'mask', 'pack']):
        return 'masky-ta-balzamy'
    if any(w in title for w in ['кондиціонер', 'conditioner', 'бальзам']):
        return 'kondytsionery'
    if any(w in title for w in ['незмивн', 'leave-in', 'leave in', 'крем для', 'cream', 'флюїд', 'fluid', 'праймер', 'primer']):
        return 'nezmyvni-zasoby'
    if any(w in title for w in ['олія', 'oil', 'сироватк', 'serum']):
        return 'oliyi-ta-syrovatky'
    if any(w in title for w in ['спрей', 'spray', 'mist', 'тонік', 'tonic']):
        return 'spreyi'

    # Fallback to parent care
    return 'doglyad-za-volossynam'


# ═══════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════
def main():
    print("=" * 60)
    print("  HAIR LAB — Catalog Data Extraction")
    print("=" * 60)

    # Parse all brands
    elgon = parse_elgon()
    mood = parse_mood()
    nevitaly = parse_nevitaly()
    inebrya = parse_inebrya()

    # Assign categories
    all_products = elgon + mood + nevitaly + inebrya
    for p in all_products:
        p['category'] = assign_category(p)

    # Generate supporting data
    categories = generate_categories()
    brands = generate_brands()
    blog_posts = generate_blog_posts()

    # Extract images
    images_manifest = extract_images()

    # Save JSON files
    def save_json(data, filename):
        path = os.path.join(SEED_DIR, filename)
        with open(path, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        print(f"  💾 Saved {filename} ({len(data)} items)")

    save_json(elgon, 'products-elgon.json')
    save_json(mood, 'products-mood.json')
    save_json(nevitaly, 'products-nevitaly.json')
    save_json(inebrya, 'products-inebrya.json')
    save_json(categories, 'categories.json')
    save_json(brands, 'brands.json')
    save_json(blog_posts, 'blog-posts.json')
    save_json(images_manifest, 'images-manifest.json')

    # Summary
    print("\n" + "=" * 60)
    print("  SUMMARY")
    print("=" * 60)
    print(f"  Elgon:    {len(elgon)} products")
    print(f"  MOOD:     {len(mood)} products")
    print(f"  Nevitaly: {len(nevitaly)} products")
    print(f"  Inebrya:  {len(inebrya)} products")
    print(f"  TOTAL:    {len(all_products)} products")
    print(f"  Images:   {len(images_manifest)} extracted")
    print(f"  Categories: {sum(1 + len(c.get('children', [])) for c in categories)}")
    print(f"  Blog posts: {len(blog_posts)}")
    print("=" * 60)

if __name__ == '__main__':
    main()
